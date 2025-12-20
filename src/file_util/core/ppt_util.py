import pptx
from io import StringIO

class PPTUtil:

    @classmethod
    def extract_texte_from_pptx(cls, filename):
        """PowerPoint(PPTX)からテキストを抽出する

        Args:
            filename: PPTXファイルパス

        Returns:
            str: 抽出されたテキスト
        """
        # 出力用のストリームを作成
        output = StringIO()
        prs = pptx.Presentation(filename)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output.write(shape.text) # type: ignore
                    output.write("\n")
        
        return output.getvalue()

